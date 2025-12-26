import sys
import base64
import os
import io
import re
import html

from typing import BinaryIO, Any
from operator import attrgetter

from ._html_converter import HtmlConverter
from ._llm_caption import llm_caption
from .._base_converter import DocumentConverter, DocumentConverterResult
from .._stream_info import StreamInfo
from .._exceptions import MissingDependencyException, MISSING_DEPENDENCY_MESSAGE
from .._conversion_quality import (
    ConversionQuality,
    FormattingLossType,
    WarningSeverity,
)

# Try loading optional (but in this case, required) dependencies
# Save reporting of any exceptions for later
_dependency_exc_info = None
try:
    import pptx
except ImportError:
    # Preserve the error and stack trace for later
    _dependency_exc_info = sys.exc_info()


ACCEPTED_MIME_TYPE_PREFIXES = [
    "application/vnd.openxmlformats-officedocument.presentationml",
]

ACCEPTED_FILE_EXTENSIONS = [".pptx"]


class PptxConverter(DocumentConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def __init__(self):
        super().__init__()
        self._html_converter = HtmlConverter()

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> bool:
        mimetype = (stream_info.mimetype or "").lower()
        extension = (stream_info.extension or "").lower()

        if extension in ACCEPTED_FILE_EXTENSIONS:
            return True

        for prefix in ACCEPTED_MIME_TYPE_PREFIXES:
            if mimetype.startswith(prefix):
                return True

        return False

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> DocumentConverterResult:
        # Check the dependencies
        if _dependency_exc_info is not None:
            raise MissingDependencyException(
                MISSING_DEPENDENCY_MESSAGE.format(
                    converter=type(self).__name__,
                    extension=".pptx",
                    feature="pptx",
                )
            ) from _dependency_exc_info[
                1
            ].with_traceback(  # type: ignore[union-attr]
                _dependency_exc_info[2]
            )

        # Perform the conversion
        presentation = pptx.Presentation(file_stream)
        md_content = ""
        slide_num = 0

        # Quality tracking
        quality = ConversionQuality(confidence=0.8)
        image_count = 0
        images_with_llm_description = 0
        images_with_alt_text = 0
        llm_caption_failures = 0
        alt_text_failures = 0
        table_count = 0
        chart_count = 0
        unsupported_chart_count = 0

        # Check if LLM is available for image descriptions
        llm_client = kwargs.get("llm_client")
        llm_model = kwargs.get("llm_model")
        llm_available = llm_client is not None and llm_model is not None
        quality.set_optional_feature("llm_image_description", llm_available)

        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title

            def get_shape_content(shape, **kwargs):
                nonlocal md_content
                nonlocal image_count, images_with_llm_description, images_with_alt_text
                nonlocal llm_caption_failures, alt_text_failures
                nonlocal table_count, chart_count, unsupported_chart_count

                # Pictures
                if self._is_picture(shape):
                    image_count += 1
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069

                    llm_description = ""
                    alt_text = ""

                    # Potentially generate a description using an LLM
                    if llm_available:
                        # Prepare a file_stream and stream_info for the image data
                        image_filename = shape.image.filename
                        image_extension = None
                        if image_filename:
                            image_extension = os.path.splitext(image_filename)[1]
                        image_stream_info = StreamInfo(
                            mimetype=shape.image.content_type,
                            extension=image_extension,
                            filename=image_filename,
                        )

                        image_stream = io.BytesIO(shape.image.blob)

                        # Caption the image
                        try:
                            llm_description = llm_caption(
                                image_stream,
                                image_stream_info,
                                client=llm_client,
                                model=llm_model,
                                prompt=kwargs.get("llm_prompt"),
                            )
                            if llm_description:
                                images_with_llm_description += 1
                        except Exception:
                            # Unable to generate a description
                            llm_caption_failures += 1

                    # Also grab any description embedded in the deck
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                        if alt_text:
                            images_with_alt_text += 1
                    except Exception:
                        # Unable to get alt text
                        alt_text_failures += 1

                    # Prepare the alt, escaping any special characters
                    alt_text = "\n".join([llm_description, alt_text]) or shape.name
                    alt_text = re.sub(r"[\r\n\[\]]", " ", alt_text)
                    alt_text = re.sub(r"\s+", " ", alt_text).strip()

                    # If keep_data_uris is True, use base64 encoding for images
                    if kwargs.get("keep_data_uris", False):
                        blob = shape.image.blob
                        content_type = shape.image.content_type or "image/png"
                        b64_string = base64.b64encode(blob).decode("utf-8")
                        md_content += f"\n![{alt_text}](data:{content_type};base64,{b64_string})\n"
                    else:
                        # A placeholder name
                        filename = re.sub(r"\W", "", shape.name) + ".jpg"
                        md_content += "\n![" + alt_text + "](" + filename + ")\n"

                # Tables
                if self._is_table(shape):
                    table_count += 1
                    md_content += self._convert_table_to_markdown(shape.table, **kwargs)

                # Charts
                if shape.has_chart:
                    chart_count += 1
                    chart_result = self._convert_chart_to_markdown(shape.chart)
                    if "[unsupported chart]" in chart_result:
                        unsupported_chart_count += 1
                    md_content += chart_result

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(
                        shape.shapes,
                        key=lambda x: (
                            float("-inf") if not x.top else x.top,
                            float("-inf") if not x.left else x.left,
                        ),
                    )
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (
                    float("-inf") if not x.top else x.top,
                    float("-inf") if not x.left else x.left,
                ),
            )
            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        # Build quality report
        quality.set_metric("slide_count", slide_num)
        quality.set_metric("image_count", image_count)
        quality.set_metric("table_count", table_count)
        quality.set_metric("chart_count", chart_count)

        # Image-related warnings
        if image_count > 0:
            if not llm_available:
                quality.add_warning(
                    f"Found {image_count} image(s) but LLM client not configured for descriptions.",
                    severity=WarningSeverity.LOW,
                    formatting_type=FormattingLossType.IMAGE_DESCRIPTION,
                    element_count=image_count,
                )
            elif llm_caption_failures > 0:
                quality.add_warning(
                    f"Failed to generate LLM descriptions for {llm_caption_failures} image(s).",
                    severity=WarningSeverity.LOW,
                    formatting_type=FormattingLossType.IMAGE_DESCRIPTION,
                    element_count=llm_caption_failures,
                )

            images_without_description = image_count - max(
                images_with_llm_description, images_with_alt_text
            )
            if images_without_description > 0:
                quality.add_warning(
                    f"{images_without_description} image(s) have no description or alt text.",
                    severity=WarningSeverity.MEDIUM,
                    formatting_type=FormattingLossType.IMAGE_DESCRIPTION,
                    element_count=images_without_description,
                )

        # Chart-related warnings
        if unsupported_chart_count > 0:
            quality.add_warning(
                f"{unsupported_chart_count} chart(s) could not be converted (unsupported type).",
                severity=WarningSeverity.MEDIUM,
                formatting_type=FormattingLossType.CHART,
                element_count=unsupported_chart_count,
            )

        # Standard PPTX formatting losses
        quality.add_formatting_loss(FormattingLossType.FONT_STYLE)
        quality.add_formatting_loss(FormattingLossType.TEXT_COLOR)

        quality.add_warning(
            "Slide transitions and animations are not preserved.",
            severity=WarningSeverity.INFO,
        )

        quality.add_warning(
            "Speaker notes are extracted but may not include rich formatting.",
            severity=WarningSeverity.INFO,
        )

        # Adjust confidence based on issues found
        if unsupported_chart_count > 0:
            quality.confidence -= 0.1 * min(unsupported_chart_count, 3)
        if image_count > 0 and not llm_available:
            quality.confidence -= 0.05

        quality.confidence = max(0.3, quality.confidence)

        return DocumentConverterResult(markdown=md_content.strip(), quality=quality)

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_table_to_markdown(self, table, **kwargs):
        # Write the table as HTML, then convert it to Markdown
        html_table = "<html><body><table>"
        first_row = True
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                if first_row:
                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                else:
                    html_table += "<td>" + html.escape(cell.text) + "</td>"
            html_table += "</tr>"
            first_row = False
        html_table += "</table></body></html>"

        return (
            self._html_converter.convert_string(html_table, **kwargs).markdown.strip()
            + "\n"
        )

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\n\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\n\n"
            data = []
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data.append(["Category"] + series_names)

            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\n".join([header, separator] + markdown_table[1:])
        except ValueError as e:
            # Handle the specific error for unsupported chart types
            if "unsupported plot type" in str(e):
                return "\n\n[unsupported chart]\n\n"
        except Exception:
            # Catch any other exceptions that might occur
            return "\n\n[unsupported chart]\n\n"
